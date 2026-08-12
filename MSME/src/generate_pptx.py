import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Brand Colors
    bg_color = RGBColor(248, 250, 252) # Dark Slate / Clinical canvas
    primary_cyan = RGBColor(2, 132, 199)
    secondary_teal = RGBColor(13, 148, 136)
    amber_color = RGBColor(234, 88, 12)
    rose_color = RGBColor(220, 38, 38)
    text_dark = RGBColor(15, 23, 42)
    muted_text = RGBColor(71, 85, 105)
    white_color = RGBColor(255, 255, 255)
    light_cyan = RGBColor(224, 242, 254)
    light_teal = RGBColor(204, 251, 241)

    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_header(slide, title_text):
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = primary_cyan

        # Accent Line
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.1), Inches(2.0), Inches(0.04))
        shape.fill.solid()
        shape.fill.fore_color.rgb = primary_cyan
        shape.line.color.rgb = primary_cyan

    # --- SLIDE 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)

    # Sub-theme badge
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(0.4), Inches(3.6), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = amber_color
    badge.line.color.rgb = amber_color
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "SUB-THEME 6: FRONTIER TECHNOLOGIES"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = white_color
    p.alignment = PP_ALIGN.CENTER

    # Logo
    if os.path.exists("MSME/images/clinical_logo.png"):
        slide1.shapes.add_picture("MSME/images/clinical_logo.png", Inches(5.8), Inches(0.8), Inches(1.6))

    # Main Title
    tb = slide1.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.733), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "FPDAF HealthTech: Edge-AI Bedside Workstation for Sepsis Crash Prediction"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = primary_cyan
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Privacy-Preserving Federated Clinical Decision Support Venture"
    p2.font.size = Pt(18)
    p2.font.color.rgb = text_dark
    p2.alignment = PP_ALIGN.CENTER

    # Card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.2), Inches(10.333), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = white_color
    card.line.color.rgb = primary_cyan
    tf = card.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "MSME Idea Hackathon 6.0 Venture Pitch | Amrita Vishwa Vidyapeetham"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = amber_color
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Lead Student Founder & Innovator: Sheela Akshar Sakhi"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = text_dark
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Faculty Mentors: Dr. Ramya G. R. (Lead Mentor) & Dr. Vandhana S. (Co-Mentor)"
    p.font.size = Pt(13)
    p.font.color.rgb = muted_text
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Department of Computer Science and Engineering, Amrita School of Computing, Coimbatore"
    p.font.size = Pt(11)
    p.font.color.rgb = secondary_teal
    p.alignment = PP_ALIGN.CENTER

    # --- SLIDE 2: What is Sepsis? ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_header(slide2, "Clinical Pathology: What is Sepsis & Why Time Matters")

    card_s1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.2))
    card_s1.fill.solid()
    card_s1.fill.fore_color.rgb = white_color
    card_s1.line.color.rgb = rose_color
    tf = card_s1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Clinical Definition of Sepsis"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = rose_color
    
    s_bullets1 = [
        "Dysregulated Immune Response: Life-threatening organ dysfunction caused by a severe systemic reaction to infection.",
        "Progression to Shock: Causes capillary leakage, tissue hypoxia, Multi-Organ Dysfunction (MODS) & circulatory crash.",
        "The Golden Window: Patient survival drops by 8% for every single hour diagnosis and treatment are delayed."
    ]
    for b in s_bullets1:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = text_dark

    card_s2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.2))
    card_s2.fill.solid()
    card_s2.fill.fore_color.rgb = white_color
    card_s2.line.color.rgb = amber_color
    tf = card_s2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Why Current ICU Alarms Fail"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = amber_color
    
    s_bullets2 = [
        "Static Threshold Alarms: Legacy monitors trigger alerts only after vitals breach rigid limits (e.g., HR > 100).",
        "Severe Alarm Fatigue: 85%-95% of bedside alarms are false positives, causing clinicians to desensitize.",
        "Missed Temporal Signals: Fails to detect complex, multi-vital co-dependencies occurring hours prior to collapse."
    ]
    for b in s_bullets2:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = text_dark

    # Bottom Impact
    scard = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.0), Inches(1.1))
    scard.fill.solid()
    scard.fill.fore_color.rgb = light_cyan
    scard.line.color.rgb = primary_cyan
    tf = scard.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Indian Healthcare Mortality Impact:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = primary_cyan
    p = tf.add_paragraph()
    p.text = "Sepsis accounts for 3.0 Million deaths annually in India (30% of total hospital deaths), making automated crash prediction an urgent national healthcare priority."
    p.font.size = Pt(12)
    p.font.color.rgb = text_dark

    # --- SLIDE 3: Market Opportunity ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_header(slide3, "Market Opportunity & Unmet Healthcare Need")

    card1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = white_color
    card1.line.color.rgb = rose_color
    tf = card1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Pain Point] High ICU Mortality Burden"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = rose_color
    
    bullets1 = [
        "3.0 Million Annual Deaths in India due to sepsis (30% of national deaths).",
        "8% Survival Drop per hour of delayed diagnosis.",
        "Traditional Alarms Fail: High false alarm rate (85%+) causes alarm fatigue."
    ]
    for b in bullets1:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = text_dark

    card2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.2))
    card2.fill.solid()
    card2.fill.fore_color.rgb = white_color
    card2.line.color.rgb = amber_color
    tf = card2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Market Gap] Cloud AI Privacy Bottleneck"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = amber_color
    
    bullets2 = [
        "DPDPA Privacy Compliance: Cloud AI startups cannot centralize hospital patient charts.",
        "Hospital Data Drift: Patient demographics vary across regions, breaking static cloud models.",
        "High Subscription Overhead: Foreign cloud solutions cost $50,000+/year."
    ]
    for b in bullets2:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = text_dark

    vcard = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.0), Inches(1.1))
    vcard.fill.solid()
    vcard.fill.fore_color.rgb = light_cyan
    vcard.line.color.rgb = primary_cyan
    tf = vcard.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "FPDAF Value Proposition:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = primary_cyan
    p = tf.add_paragraph()
    p.text = "FPDAF delivers a zero-cloud-upload Edge-AI workstation node that trains collaboratively across hospitals while guaranteeing 100% DPDPA data privacy."
    p.font.size = Pt(12)
    p.font.color.rgb = text_dark

    # --- SLIDE 4: Architecture Diagram ---
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_header(slide4, "Product Architecture & Edge Node Dataflow")

    if os.path.exists("MSME/images/fpdafarcdiagram.png"):
        slide4.shapes.add_picture("MSME/images/fpdafarcdiagram.png", Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.0))

    # --- SLIDE 5: Proprietary IP ---
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_header(slide5, "Proprietary IP & Innovation Moat")

    ips = [
        ("1. Dynamic Dual-Phase (DDP) Attention", "2D cross-attention tracking 24h sequence hours AND multi-modal vital co-dependencies (e.g., HR vs Blood Pressure).", primary_cyan),
        ("2. AD-CUSUM Drift Auditing", "Real-time loss residual control charts flagging patient population shifts on edge nodes before model degradation occurs.", secondary_teal),
        ("3. Client-Side Selective Personalization", "Freezes shared LSTM backbones upon drift alerts, fine-tuning only local heads (38% bandwidth savings, 6-min adaptation).", amber_color),
        ("4. Temporal XAI Saliency Heatmaps", "Generates hourly bedside vital contribution heatmaps for ICU doctor verification.", rose_color)
    ]

    coords = [(0.6, 1.5), (6.8, 1.5), (0.6, 4.3), (6.8, 4.3)]
    for idx, (title, desc, col) in enumerate(ips):
        x, y = coords[idx]
        c = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2.4))
        c.fill.solid()
        c.fill.fore_color.rgb = white_color
        c.line.color.rgb = col
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = text_dark

    # --- SLIDE 6: Competitor Table ---
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_header(slide6, "Competitor Analysis & Cost Benchmarking")

    table_shape = slide6.shapes.add_table(5, 5, Inches(0.6), Inches(1.4), Inches(12.133), Inches(3.2))
    table = table_shape.table

    headers = ["Competitor System", "Offering / Model", "Privacy & Drift Limit", "Deployment Cost", "Edge-AI"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary_cyan
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = white_color

    data = [
        ["Epic Deterioration Index", "Cloud SaaS EHR scoring", "Cloud-only (Violates DPDPA)", "$50,000+/yr (~Rs. 40L)", "No"],
        ["Philips IntelliVue", "Bedside hardware monitor", "Threshold alarms (No AI)", "Rs. 15-20 Lakhs/bed", "No"],
        ["GE Healthcare CARESCAPE", "Enterprise server software", "Vendor lock-in, No privacy", "Rs. 25-30 Lakhs/setup", "No"],
        ["FPDAF HealthTech (Ours)", "Federated Edge-AI Node", "100% DPDPA, Real-time Drift", "Rs. 1.5 - 2.0 Lakhs/node", "YES"]
    ]

    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if row_idx == 3:
                cell.fill.fore_color.rgb = light_teal
            else:
                cell.fill.fore_color.rgb = white_color
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(11)
            if row_idx == 3:
                p.font.bold = True
                p.font.color.rgb = secondary_teal
            else:
                p.font.color.rgb = text_dark

    c_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.9), Inches(12.133), Inches(2.1))
    c_card.fill.solid()
    c_card.fill.fore_color.rgb = white_color
    c_card.line.color.rgb = secondary_teal
    tf = c_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Commercial Disruption & Advantage for Indian Hospitals:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = secondary_teal
    
    p = tf.add_paragraph()
    p.text = "• 85%-90% Capital Cost Savings: Rs. 1.5L - 2.0L per node vs Rs. 15L-40L for global legacy systems."
    p.font.size = Pt(12)
    p.font.color.rgb = text_dark
    p = tf.add_paragraph()
    p.text = "• 100% DPDPA 2023 Privacy Compliance: Zero cloud patient chart uploads, eliminating legal liabilities."
    p.font.size = Pt(12)
    p.font.color.rgb = text_dark

    # --- SLIDE 7: Operational Video Demo (EMBEDDED VIDEO) ---
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)
    add_header(slide7, "Venture Validation & Operational Video Demo (TRL 4/5)")

    card_l = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.4))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = white_color
    card_l.line.color.rgb = primary_cyan
    tf = card_l.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Product Benchmark Performance"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = primary_cyan
    
    bmark_bullets = [
        "86.51% Accuracy: Superior to FedAvg (75.94%) and FedProx (78.85%).",
        "0.8214 AUROC Score: High predictive accuracy for early crash detection.",
        "38% Bandwidth Savings: 1.52 GB vs 2.48 GB in standard FL.",
        "Operational Tech Stack: FastAPI REST backend, SQLite warehouse (`clinical_warehouse.db`), and React Vite UI."
    ]
    for b in bmark_bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = text_dark

    v_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.4))
    v_card.fill.solid()
    v_card.fill.fore_color.rgb = white_color
    v_card.line.color.rgb = amber_color
    tf = v_card.text_frame
    p = tf.paragraphs[0]
    p.text = "Bedside Workstation Video Demo"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = amber_color

    video_path = "assets/FPDAF.mp4" if os.path.exists("assets/FPDAF.mp4") else "FPDAF.mp4"
    poster_path = "MSME/images/dashboarddoctor.png"
    if os.path.exists(video_path):
        movie_shape = slide7.shapes.add_movie(
            video_path,
            Inches(7.1), Inches(2.2), Inches(5.2), Inches(3.4),
            poster_frame_image=poster_path if os.path.exists(poster_path) else None,
            mime_type="video/mp4"
        )

    cap_tb = slide7.shapes.add_textbox(Inches(6.8), Inches(5.8), Inches(5.8), Inches(1.0))
    tf = cap_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "▶ Click to Play Operational Workstation Video Demo"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = amber_color
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Real-time telemetry, AD-CUSUM drift curves & XAI heatmaps."
    p2.font.size = Pt(11)
    p2.font.color.rgb = muted_text
    p2.alignment = PP_ALIGN.CENTER

    # --- SLIDE 8: Financial Breakdown ---
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8)
    add_header(slide8, "Itemized Technology Expenditure Breakdown (Rs. 10.0 Lakhs)")

    fin_table_shape = slide8.shapes.add_table(4, 3, Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.4))
    fin_table = fin_table_shape.table

    f_headers = ["Component Category", "Itemized Specs & Market Research Quote", "Cost Outlay"]
    for i, h in enumerate(f_headers):
        cell = fin_table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary_cyan
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = white_color

    fin_data = [
        ["A. Edge Hardware Prototyping & MSME Assembly", "• 3x NVIDIA Jetson Orin Nano (8GB) Kits @ Rs 48K = 1.44L\n• 3x Medical 15.6\" Touch Monitors & HDMI Controllers = 0.66L\n• 3x NEMA Metallic Bedside Cabinets (Coimbatore MSME) = 0.75L\n• Custom SMT Carrier PCB & Medical Power Assembly = 0.75L", "Rs. 3.60 Lakhs"],
        ["B. EHR Data Curation & Clinical Annotation", "• PhysioNet ICU Cohort Cleaning (40,000 vital records) = 1.40L\n• Senior ICU Specialist Onset Audit (70h @ Rs 2,500/h) = 1.75L\n• Vital Stream Generator Calibration (scaler.pkl) = 0.25L", "Rs. 3.40 Lakhs"],
        ["C. Cybersecurity & Compliance Audits", "• CERT-In External Vulnerability & Penetration Audit = 1.80L\n• DPDPA 2023 Data Privacy Encryption Compliance Audit = 1.20L", "Rs. 3.00 Lakhs"]
    ]

    for r_idx, r_data in enumerate(fin_data):
        for c_idx, text in enumerate(r_data):
            cell = fin_table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = white_color
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(11)
            if c_idx == 0:
                p.font.bold = True
                p.font.color.rgb = primary_cyan if r_idx == 0 else (secondary_teal if r_idx == 1 else amber_color)
            elif c_idx == 2:
                p.font.bold = True
                p.font.color.rgb = primary_cyan if r_idx == 0 else (secondary_teal if r_idx == 1 else amber_color)
            else:
                p.font.color.rgb = text_dark

    # --- SLIDE 9: Business Model ---
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9)
    add_header(slide9, "Business Model & Local MSME Partnership")

    b_card1 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.3))
    b_card1.fill.solid()
    b_card1.fill.fore_color.rgb = white_color
    b_card1.line.color.rgb = primary_cyan
    tf = b_card1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Commercialization & Revenue Streams"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = primary_cyan
    
    b_bullets1 = [
        "1. B2B Edge Node Appliance: Direct sales of turnkey Jetson hardware workstations at Rs. 1.50L -- 2.00L / node.",
        "2. SaaS Telemetry Subscription: Annual maintenance, model updates & telemetry auditing at Rs. 25,000 / node / year.",
        "3. OEM Engine Licensing: Licensing PyTorch FL engine to monitor OEMs (Philips, Mindray)."
    ]
    for b in b_bullets1:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = text_dark

    b_card2 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.3))
    b_card2.fill.solid()
    b_card2.fill.fore_color.rgb = white_color
    b_card2.line.color.rgb = secondary_teal
    tf = b_card2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Coimbatore MSME Hardware Ecosystem"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = secondary_teal
    
    b_bullets2 = [
        "PCB Electronics Assembly: SMT component mounting & carrier board fabrication with Coimbatore electronics MSMEs.",
        "Metallic Cabinet Enclosure: Powder-coated sheet metal node fabrication with local industrial MSMEs.",
        "12-Month Target: ISO 13485 quality pack, CDSCO Class B filing & 3-hospital tele-ICU pilot run."
    ]
    for b in b_bullets2:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = text_dark

    som_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.0), Inches(12.0), Inches(1.0))
    som_card.fill.solid()
    som_card.fill.fore_color.rgb = light_cyan
    som_card.line.color.rgb = amber_color
    tf = som_card.text_frame
    p = tf.paragraphs[0]
    p.text = "Market Scaling Traction:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = amber_color
    p = tf.add_paragraph()
    p.text = "Target Year 1--2 SOM of Rs. 40 Crores across 12 regional hospital chains representing 2,500 ICU beds."
    p.font.size = Pt(12)
    p.font.color.rgb = text_dark

    # --- SLIDE 10: Conclusion ---
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10)

    tb = slide10.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.733), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = primary_cyan
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "FPDAF HealthTech: Building India's Next-Gen Privacy-Preserving MedTech AI"
    p2.font.size = Pt(20)
    p2.font.color.rgb = text_dark
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Questions & Investor Discussion"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = secondary_teal
    p3.alignment = PP_ALIGN.CENTER

    output_path = "MSME/msmepresentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated 10-slide PowerPoint deck with embedded video at: {output_path}")

if __name__ == "__main__":
    create_deck()
